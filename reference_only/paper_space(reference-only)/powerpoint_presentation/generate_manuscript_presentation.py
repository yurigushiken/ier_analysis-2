"""
Generate a PI-facing PowerPoint presentation from the paper_space manuscript.

Usage (from repo root):
    python paper_space/powerpoint_presentation/generate_manuscript_presentation.py

Outputs into the same directory:
    pi_boundary_manuscript_presentation.pptx
    pi_boundary_manuscript_presentation_qc.json
    slide_previews/   (one JPEG per slide for visual QA)
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
REPO_ROOT = Path(__file__).resolve().parents[2]
OUTPUT_DIR = Path(__file__).resolve().parent
OUTPUT_PPTX = OUTPUT_DIR / "pi_boundary_manuscript_presentation.pptx"
OUTPUT_QC = OUTPUT_DIR / "pi_boundary_manuscript_presentation_qc.json"
PREVIEW_DIR = OUTPUT_DIR / "slide_previews"

# ---------------------------------------------------------------------------
# Slide dimensions (16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Color palette: subtle beige + light blue
# ---------------------------------------------------------------------------
COLOR_BG = RGBColor(245, 240, 232)
COLOR_ACCENT = RGBColor(156, 196, 228)  # stronger light blue for accents
COLOR_ACCENT_LIGHT = RGBColor(210, 228, 243)  # very light blue tint
COLOR_TEXT = RGBColor(38, 50, 66)
COLOR_MUTED = RGBColor(90, 110, 130)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_PLACEHOLDER = RGBColor(230, 235, 240)
COLOR_DARK_ACCENT = RGBColor(62, 100, 140)

# ---------------------------------------------------------------------------
# Typography
# ---------------------------------------------------------------------------
FONT_TITLE = Pt(38)
FONT_SUBTITLE = Pt(22)
FONT_HEADING = Pt(30)
FONT_BODY = Pt(22)
FONT_CAPTION = Pt(14)
FONT_SMALL = Pt(11)
FONT_FAMILY = "Calibri"


# ===== Layout helpers ======================================================

def _set_font(run, *, size=FONT_BODY, bold=False, italic=False,
              color=COLOR_TEXT, name=FONT_FAMILY):
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = name


def add_bg(slide):
    """Full-slide beige background."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    r.fill.solid()
    r.fill.fore_color.rgb = COLOR_BG
    r.line.fill.background()


def add_header_bar(slide):
    """Thin accent strip at the very top of the slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()


def add_footer(slide, text: str = "CCN 2026 | PI/ANS Boundary at 4-5"):
    """Small footer line at slide bottom."""
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    _set_font(run, size=FONT_SMALL, color=COLOR_MUTED)


def add_heading(slide, text: str, *, y=Inches(0.22)):
    """Slide heading with underline accent bar."""
    box = slide.shapes.add_textbox(Inches(0.7), y, Inches(11.9), Inches(0.85))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, size=FONT_HEADING, bold=True, color=COLOR_TEXT)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), y + Inches(0.82),
        Inches(3.0), Inches(0.055))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()


def add_slide_number(slide, num: int):
    box = slide.shapes.add_textbox(
        Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(num)
    _set_font(run, size=FONT_SMALL, color=COLOR_MUTED)


def _new_slide(prs, num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide)
    add_footer(slide)
    add_slide_number(slide, num)
    return slide


# ===== Figure placement (aspect-ratio aware) ===============================

def _img_dims(path: Path):
    """Return (w_px, h_px) or None."""
    if path and path.exists():
        from PIL import Image
        im = Image.open(path)
        return im.width, im.height
    return None


def add_image_fit(slide, img_path: Path | None, left, top, max_w, max_h,
                  qc: dict):
    """Insert image preserving aspect ratio inside a bounding box."""
    if img_path is None:
        return False
    if not img_path.exists():
        _add_placeholder(slide, left, top, max_w, max_h)
        qc["missing_assets"].append(str(img_path))
        return False

    dims = _img_dims(img_path)
    if dims:
        w_px, h_px = dims
        ratio = w_px / h_px
        # Fit within max_w x max_h
        box_ratio = Emu(max_w) / Emu(max_h)
        if ratio >= box_ratio:
            # Width-limited
            final_w = max_w
            final_h = int(Emu(max_w) / ratio)
        else:
            # Height-limited
            final_h = max_h
            final_w = int(Emu(max_h) * ratio)
        # Center within bounding box
        x_offset = (Emu(max_w) - final_w) // 2
        y_offset = (Emu(max_h) - final_h) // 2
        slide.shapes.add_picture(
            str(img_path),
            Emu(Emu(left) + x_offset), Emu(Emu(top) + y_offset),
            width=final_w, height=final_h)
    else:
        slide.shapes.add_picture(str(img_path), left, top, width=max_w)

    qc["included_assets"].append(str(img_path))
    return True


def _add_placeholder(slide, left, top, w, h):
    ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    ph.fill.solid()
    ph.fill.fore_color.rgb = COLOR_PLACEHOLDER
    ph.line.color.rgb = COLOR_ACCENT
    tf = ph.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "[Figure not found]"
    _set_font(run, size=Pt(14), color=COLOR_MUTED, italic=True)
    p.alignment = PP_ALIGN.CENTER


# ===== Bullet text helper ==================================================

def add_bullets(slide, bullets: list[str | tuple[str, int]],
                *, left=Inches(0.9), top=Inches(1.45),
                width=Inches(11.5), height=Inches(5.5)):
    """Add bullet list. Each entry is a string or (string, indent_level)."""
    BULLET_CHAR = "*  "  # bullet
    SUB_BULLET = "-  "  # hyphen for sub-items
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()
    for idx, item in enumerate(bullets):
        text, level = (item, 0) if isinstance(item, str) else item
        prefix = SUB_BULLET if level > 0 else BULLET_CHAR
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = prefix + text
        _set_font(run, size=FONT_BODY if level == 0 else Pt(19),
                  color=COLOR_TEXT if level == 0 else COLOR_MUTED)
        p.space_after = Pt(16) if level == 0 else Pt(10)
        p.space_before = Pt(4) if level > 0 else Pt(0)
        if level > 0:
            p.margin_left = Inches(0.4)


def add_text_panel(slide, lines: list[str], left, top, width, height,
                   *, bg=COLOR_ACCENT_LIGHT):
    """Text box with light tinted background for right-side narrative."""
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = bg
    panel.line.fill.background()
    tf = panel.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        _set_font(run, size=Pt(18), color=COLOR_TEXT)
        p.space_after = Pt(12)


def add_caption(slide, text: str, left, top, width):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, size=FONT_CAPTION, color=COLOR_MUTED, italic=True)


# ===== Asset registry ======================================================

def build_assets() -> dict[str, Path]:
    R = REPO_ROOT
    return {
        "trial_structure": R / "paper_space/media/trial_structure.png",
        "conditions_24": R / "paper_space/media/conditions-24.png",
        "conditions_6": R / "paper_space/media/conditions-6.png",
        "distance_with1": R / "results/analysis-lda/temporal_change24_all-100-700ms/figures/temporal_distance_timecourse.png",
        "distance_no1": R / "results/analysis-lda/temporal_change24_all-100-700ms/figures-no1/temporal_distance_timecourse.png",
        "static_c6_rdms": R / "results/analysis-lda/static-change6/static_change6_all/figures-pairwise/temporal_model_rdms_all_compact_row-static_change6_all-pairwise.png",
        "static_c5_rdms": R / "results/analysis-lda/static-change5/static_change5_all/figures-pairwise/temporal_model_rdms_all_compact_row-static_change5_all-pairwise.png",
        "static_mds": R / "results/analysis-lda/static-change6/static_change6_all/figures-pairwise/mds_category-static_change6_all-pairwise-notitle.png",
        "fits_with1": R / "results/analysis-lda/temporal_change24_all-100-700ms/figures/temporal_model_fits_with_significance-temporal_change24_all-100-700ms-with_RT.png",
        "fits_no1": R / "results/analysis-lda/temporal_change24_all-100-700ms/figures-no1/temporal_model_fits_with_significance-temporal_change24_all-100-700ms-no1-with_RT.png",
        "snapshots": R / "results/analysis-lda/temporal_change24_all-100-700ms/figures/temporal_rdm_snapshots-temporal_change24_all-100-700ms.png",
        "compact_with1": R / "results/analysis-lda/temporal_change24_all-100-700ms/figures/temporal_model_rdms_all_compact_row-temporal_change24_all-100-700ms-RT.png",
        "compact_no1": R / "results/analysis-lda/temporal_change24_all-100-700ms/figures-no1/temporal_model_rdms_all_compact_row-temporal_change24_all-100-700ms-no1-RT.png",
        "temporal_mds": R / "results/analysis-lda/temporal_change24_all-100-700ms/figures/temporal_mds_static_0-700ms-change_code-temporal_change24_all-100-700ms.png",
        "workflow": OUTPUT_DIR / "media" / "workflow_diagram.png",
        "sliding_window": OUTPUT_DIR / "media" / "sliding_window_diagram.png",
    }


# ===== Slide builders (one per slide) ======================================

def slide_01_title(prs, n, _a, _q):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    # Top accent bar wider
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    # Title
    box = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.6), Inches(11.3), Inches(2.4))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = ("Neural Similarity Patterns Support a Categorical "
                "Boundary Between 4 and 5 in Numerosity (1-6)")
    _set_font(run, size=FONT_TITLE, bold=True, color=COLOR_TEXT)

    # Subtitle
    box2 = slide.shapes.add_textbox(
        Inches(1.0), Inches(4.3), Inches(11.3), Inches(0.6))
    p2 = box2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "CCN 2026  *  PI Discussion"
    _set_font(run2, size=FONT_SUBTITLE, color=COLOR_ACCENT)

    # Meta
    box3 = slide.shapes.add_textbox(
        Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.0))
    tf3 = box3.text_frame
    for txt in ["Dataset: ALL trials (correct + incorrect), 24 participants",
                "EEG oddball change-detection task with dot arrays (1-6)"]:
        p3 = tf3.paragraphs[0] if txt == tf3.text else tf3.add_paragraph()
        run3 = p3.add_run()
        run3.text = txt
        _set_font(run3, size=Pt(18), color=COLOR_MUTED)
        p3.space_after = Pt(6)

    add_slide_number(slide, n)


def slide_02_question(prs, n, _a, _q):
    slide = _new_slide(prs, n)
    add_heading(slide, "Research Questions and Hypotheses")
    add_bullets(slide, [
        "RQ1: Does neural representational similarity support a categorical "
        "boundary between PI and ANS number systems?",
        "RQ2: Does the boundary fall at 4-5 rather than the traditionally "
        "assumed 3-4?",
        "RQ3: Does the preceding prime stimulus influence target "
        "representational geometry?",
        "H1: Boundary-based PI models (at 4-5) will outperform both the "
        "3-4 alternative and graded ANS log-ratio models.",
        "H2: Excluding numerosity 1 will attenuate ANS fits while preserving "
        "the 4-5 boundary signal.",
        "H3: Prime x Target PI models will show sustained significance, "
        "reflecting latent working-memory contributions.",
    ])


def slide_03_numerosity1(prs, n, _a, _q):
    slide = _new_slide(prs, n)
    add_heading(slide, 'Why Numerosity 1 Matters: the "No-1" Control')
    add_bullets(slide, [
        "Under Weber-law scaling, log(n/1) is the largest possible "
        "log-ratio for any pair - numerosity 1 is maximally distant.",
        "This inflates ANS model fits, making it hard to tell if strong "
        "ANS performance reflects genuine ratio sensitivity or just "
        "the discriminability of 1.",
        "Solution: repeat all analyses after excluding pairs containing "
        'numerosity 1 (the "no-1" control).',
        "Goal: separate true boundary structure from singleton salience.",
    ])


def slide_04_rsa(prs, n, _a, _q):
    slide = _new_slide(prs, n)
    add_heading(slide, "RSA Strategy and Two-Stage Pipeline")
    add_bullets(slide, [
        "Single-trial EEG -> pairwise LDA decoding -> neural RDMs.",
        ("Stage 1: Whole-epoch target RSA (0-700 ms)", 0),
        ("Collapses across prime identity; 6x6 target RDM; maximizes SNR.", 1),
        ("Stage 2: Time-resolved 24-condition prime-target RSA", 0),
        ("Sliding window (-100 to 700 ms); full 24x24 RDM preserving "
         "prime and target identity.", 1),
        "Spearman correlations + cluster-based permutation inference.",
    ])


def slide_05_models(prs, n, _a, _q):
    slide = _new_slide(prs, n)
    add_heading(slide, "Model Family Overview")
    add_bullets(slide, [
        "Target PI (boundary at 3-4 or 4-5): categorical model. "
        "Within PI: uniform dissimilarity. Within ANS: log-ratio. "
        "Cross-boundary: elevated baseline + log-ratio.",
        "Prime x Target PI: encodes PI/ANS membership of both "
        "preceding prime and current target (0, 1, or 2 mismatches).",
        "ANS (log-ratio): pure Weber-law distance, no categorical boundary.",
        "Pixel: total white-pixel area difference (low-level visual control).",
        "RT: absolute reaction-time difference (task-difficulty control).",
    ])


def slide_06_task(prs, n, a, q):
    slide = _new_slide(prs, n)
    add_heading(slide, "Task Design and Conditions")
    # Trial structure: very wide image -> full width
    add_image_fit(slide, a["trial_structure"],
                  Inches(0.6), Inches(1.2), Inches(12.1), Inches(2.5), q)
    add_caption(slide, "Single-trial structure: 3-5 primes -> target",
                Inches(0.7), Inches(3.75), Inches(5.0))
    # Conditions 24 + conditions 6 side by side
    add_image_fit(slide, a["conditions_24"],
                  Inches(0.6), Inches(4.15), Inches(6.8), Inches(3.0), q)
    add_caption(slide, "24 prime-target conditions",
                Inches(0.7), Inches(6.85), Inches(4.0))
    add_image_fit(slide, a["conditions_6"],
                  Inches(7.6), Inches(4.5), Inches(5.2), Inches(1.5), q)
    add_caption(slide, "6 target numerosities (collapsed)",
                Inches(7.7), Inches(6.1), Inches(4.0))


def slide_07_methods(prs, n, _a, _q):
    slide = _new_slide(prs, n)
    add_heading(slide, "Methods Summary")
    add_bullets(slide, [
        "24 adults, oddball change detection, dot arrays 1-6.",
        "128-channel EEG, HAPPE preprocessing (1.5-35 Hz, "
        "wavelet artifact correction).",
        "LDA with Ledoit-Wolf shrinkage, GroupKFold CV (6 folds, "
        "subject-level independence).",
        ("Whole-epoch: one Spearman rho per subject per model; "
         "one-sample t-test with Holm correction.", 0),
        ("Temporal: cluster permutation (5,000 permutations); "
         "model-difference contrasts (10,000 permutations, Fisher-z).", 0),
    ])


def slide_08_distance(prs, n, a, q):
    slide = _new_slide(prs, n)
    add_heading(slide, "Pairwise Decoding by Numerical Distance")
    # Two landscape panels side by side
    add_image_fit(slide, a["distance_with1"],
                  Inches(0.5), Inches(1.25), Inches(6.2), Inches(4.2), q)
    add_caption(slide, "(a) All pairs: clear monotonic distance gradient",
                Inches(0.6), Inches(5.5), Inches(6.0))
    add_image_fit(slide, a["distance_no1"],
                  Inches(6.8), Inches(1.25), Inches(6.2), Inches(4.2), q)
    add_caption(slide, "(b) No-1: gradient largely collapses",
                Inches(6.9), Inches(5.5), Inches(6.0))
    # Key takeaway bar
    add_text_panel(slide,
                   ["Numerosity 1 disproportionately drives pairwise "
                    "discriminability and the distance gradient."],
                   Inches(0.5), Inches(6.1), Inches(12.5), Inches(0.7))


def slide_09_table1(prs, n, _a, _q):
    slide = _new_slide(prs, n)
    add_heading(slide, "Whole-Epoch Model Comparison")
    # Recreate Table 1 as a native PPTX table
    rows, cols = 8, 5
    tbl_shape = slide.shapes.add_table(
        rows, cols, Inches(0.7), Inches(1.25), Inches(11.9), Inches(4.5))
    tbl = tbl_shape.table
    col_widths = [Inches(3.0), Inches(1.5), Inches(2.0), Inches(1.5), Inches(2.0)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    data = [
        ["Model", "r (1-6)", "p (1-6)", "r (2-6)", "p (2-6)"],
        ["Target PI(x-4)", ".289", "< .001", ".269", ".003"],
        ["Target PI(x-3)", ".207", "< .001", ".076", "ns"],
        ["ANS", ".313", "< .001", ".074", "ns"],
        ["Pixel", ".056", "ns", ".069", "ns"],
        ["RT", ".162", ".005", ".268", ".007"],
        ["Noise ceiling", ".417", "", ".332", ""],
        ["Boundary PI(x-4) vs PI(x-3)", "d = 0.53", "p = .017",
         "d = 0.79", "p < .001"],
    ]
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    _set_font(run, size=Pt(17),
                              bold=(r_idx == 0 or r_idx == 7),
                              color=COLOR_TEXT if r_idx > 0 else COLOR_WHITE)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_DARK_ACCENT
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_ACCENT_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE

    add_text_panel(slide, [
        "Key: In no-1 (2-6), only Target PI(2-4) remains "
        "significant. Boundary contrast strengthens to d = 0.79."
    ], Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.7))


def slide_10_static_rdms(prs, n, a, q):
    """Static RDMs: these are very wide (6.62:1) -- stack vertically."""
    slide = _new_slide(prs, n)
    add_heading(slide, "Whole-Epoch Brain and Model RDMs")
    add_image_fit(slide, a["static_c6_rdms"],
                  Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.5), q)
    add_caption(slide, "(a) Target-only (1-6): Brain RDM + model RDMs",
                Inches(0.6), Inches(3.7), Inches(8.0))
    add_image_fit(slide, a["static_c5_rdms"],
                  Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.5), q)
    add_caption(slide, "(b) No-1 (2-6): ANS gradient attenuated; "
                "PI(2-4) structure persists",
                Inches(0.6), Inches(6.6), Inches(10.0))


def slide_11_static_mds(prs, n, a, q):
    slide = _new_slide(prs, n)
    add_heading(slide, "Whole-Epoch MDS: Representational Geometry")
    add_image_fit(slide, a["static_mds"],
                  Inches(0.6), Inches(1.15), Inches(6.5), Inches(5.5), q)
    add_text_panel(slide, [
        "Numerosity 1 (orange) is isolated from the rest.",
        "PI range (2-4, blue) clusters together.",
        "ANS range (5-6, green) occupies a separate region.",
        "Consistent with a categorical boundary at 4-5.",
    ], Inches(7.5), Inches(1.5), Inches(5.2), Inches(4.0))


def slide_12_temporal_fits(prs, n, a, q):
    slide = _new_slide(prs, n)
    add_heading(slide, "Time-Resolved Model Fits")
    add_image_fit(slide, a["fits_with1"],
                  Inches(0.4), Inches(1.15), Inches(6.3), Inches(4.5), q)
    add_caption(slide,
                "(a) With-1: Prime x Target PI(1-4) "
                "longest cluster, 70-650 ms",
                Inches(0.5), Inches(5.7), Inches(6.1))
    add_image_fit(slide, a["fits_no1"],
                  Inches(6.8), Inches(1.15), Inches(6.3), Inches(4.5), q)
    add_caption(slide,
                "(b) No-1: Prime x Target PI(2-4) "
                "retains broad cluster, 90-630 ms",
                Inches(6.9), Inches(5.7), Inches(6.1))
    add_text_panel(slide, [
        "Uninterrupted 580 ms cluster from Prime x Target PI(1-4) "
        "is the longest among all tested models."
    ], Inches(0.4), Inches(6.35), Inches(12.5), Inches(0.6))


def slide_13_compact_rdms(prs, n, a, q):
    """Wide row-format compact RDMs (7.75:1) -- stack vertically like slide 10."""
    slide = _new_slide(prs, n)
    add_heading(slide, "24-Condition Brain and Model RDMs")
    # Stacked vertically, full width
    add_image_fit(slide, a["compact_with1"],
                  Inches(0.4), Inches(1.25), Inches(12.5), Inches(2.4), q)
    add_caption(slide,
                "(a) All 24 conditions: Brain RDM + model RDMs (with RT)",
                Inches(0.5), Inches(3.65), Inches(10.0))
    add_image_fit(slide, a["compact_no1"],
                  Inches(0.4), Inches(4.0), Inches(12.5), Inches(2.4), q)
    add_caption(slide,
                "(b) No-1 subset: PI(2-4) block structure persists; "
                "ANS gradient attenuated",
                Inches(0.5), Inches(6.4), Inches(10.0))


def slide_14_snapshots(prs, n, a, q):
    """RDM evolution: very wide image (2:1) -- full width."""
    slide = _new_slide(prs, n)
    add_heading(slide, "RDM Evolution Over Time")
    add_image_fit(slide, a["snapshots"],
                  Inches(0.3), Inches(1.15), Inches(12.7), Inches(4.3), q)
    add_caption(slide,
                "100 ms bins from -100 to 700 ms. Categorical "
                "structure emerges by ~70 ms and persists.",
                Inches(0.4), Inches(5.5), Inches(12.0))
    add_text_panel(slide, [
        "Early bins (-100-0 ms): no structure.  "
        "100-200 ms: strong target-1 differentiation.  "
        "Maintained through 600-700 ms.",
    ], Inches(0.4), Inches(5.95), Inches(12.5), Inches(0.6))


def slide_15_model_diff(prs, n, _a, _q):
    slide = _new_slide(prs, n)
    add_heading(slide, "Model-Difference Tests")
    rows, cols = 8, 4
    tbl_shape = slide.shapes.add_table(
        rows, cols, Inches(0.7), Inches(1.25), Inches(11.9), Inches(4.0))
    tbl = tbl_shape.table
    widths = [Inches(4.5), Inches(2.5), Inches(2.5), Inches(1.5)]
    for i, w in enumerate(widths):
        tbl.columns[i].width = w

    data = [
        ["Comparison", "Best fit", "Range (ms)", "p"],
        ["With-1: PI(1-4) vs PI(1-3)", "PI(1-4)",
         "230-330", "< .001"],
        ["With-1: PI(1-4) vs PI(1-3)", "PI(1-4)",
         "490-550", ".039"],
        ["With-1: PI(1-4) vs PI(1-3)", "PI(1-4)",
         "570-650", ".014"],
        ["With-1: PI(1-4) vs ANS", "ANS", "90-150", ".050"],
        ["No-1: PI(2-4) vs ANS", "PI(2-4)",
         "170-310", "< .001"],
        ["No-1: PI(2-4) vs PI(2-3)", "PI(2-4)",
         "130-350", "< .001"],
        ["No-1: PI(2-4) vs PI(2-3)", "PI(2-4)",
         "390-550", "< .001"],
    ]
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    _set_font(run, size=Pt(16),
                              bold=(r_idx == 0),
                              color=COLOR_TEXT if r_idx > 0 else COLOR_WHITE)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_DARK_ACCENT
            elif r_idx == 4:
                # highlight the ANS row differently
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(252, 240, 228)
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_ACCENT_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE

    add_text_panel(slide, [
        "4-5 boundary consistently outperforms 3-4 across "
        "multiple time windows. ANS advantage is limited to "
        "90-150 ms with-1 (borderline p = .050).",
    ], Inches(0.7), Inches(5.45), Inches(11.9), Inches(0.85))


def slide_16_temporal_mds(prs, n, a, q):
    slide = _new_slide(prs, n)
    add_heading(slide, "24-Condition MDS: Representational Geometry")
    add_image_fit(slide, a["temporal_mds"],
                  Inches(0.5), Inches(1.15), Inches(7.0), Inches(5.5), q)
    add_text_panel(slide, [
        "Each dot is a prime-target condition (e.g., 42 = prime 4, "
        "target 2).",
        "Yellow = target 1 (isolated cluster).",
        "Blue = PI targets 2-4 (central group).",
        "Green = ANS targets 5-6 (separate region).",
        "Separation by target system membership is the dominant "
        "dimension of representational structure.",
    ], Inches(7.8), Inches(1.4), Inches(5.0), Inches(5.0))


def slide_17_discussion(prs, n, _a, _q):
    slide = _new_slide(prs, n)
    add_heading(slide, "Discussion")
    add_bullets(slide, [
        "Converging analyses place the PI/ANS boundary at 4-5, "
        "not 3-4.",
        "Prime information contributes: Prime x Target PI(1-4) "
        "achieves the longest contiguous cluster (70-650 ms), "
        "consistent with activity-silent working memory.",
        'Numerosity 1 is "special": inflates ANS fits, '
        "but removing it preserves the 4-5 boundary signal.",
        "RT and boundary structure overlap (r ~ .27 in no-1), "
        "but boundary-location contrast dissociates the two "
        "(PI(2-4) > PI(2-3)).",
    ])


def slide_19_workflow(prs, n, a, q):
    slide = _new_slide(prs, n)
    add_heading(slide, "Analysis Pipeline (Supplementary)")
    add_image_fit(slide, a["workflow"],
                  Inches(0.3), Inches(1.15), Inches(12.7), Inches(5.8), q)


def slide_20_sliding_window(prs, n, a, q):
    slide = _new_slide(prs, n)
    add_heading(slide, "Sliding Window Approach (Supplementary)")
    add_image_fit(slide, a["sliding_window"],
                  Inches(0.3), Inches(1.15), Inches(12.7), Inches(5.8), q)


def slide_18_limitations(prs, n, _a, _q):
    slide = _new_slide(prs, n)
    add_heading(slide, "Limitations, Conclusion, and Next Steps")
    add_bullets(slide, [
        "Dot arrays covary with area/density; Pixel control is "
        "non-significant but visual confounds cannot be fully ruled out.",
        "ANS range contains only 5-6 - limits within-ANS "
        "granularity tests.",
        "PI modeled as categorical; within-PI graded structure not tested.",
        "Findings are paradigm-specific; replication needed.",
    ], top=Inches(1.25), height=Inches(3.2))
    # Conclusion box
    add_text_panel(slide, [
        "Conclusion: PI/ANS boundary in this dataset is best captured "
        "at 4-5. Both whole-epoch and temporal RSA converge on this "
        "result.",
        "Next steps: replicate in alternative paradigms, expand ANS "
        "range, and test within-PI graded models.",
    ], Inches(0.7), Inches(4.8), Inches(11.9), Inches(2.0),
                   bg=COLOR_ACCENT_LIGHT)


# ===== Master build ========================================================

SLIDE_BUILDERS = [
    slide_01_title,
    slide_02_question,
    slide_03_numerosity1,
    slide_04_rsa,
    slide_05_models,
    slide_06_task,
    slide_07_methods,
    slide_08_distance,
    slide_09_table1,
    slide_10_static_rdms,
    slide_11_static_mds,
    slide_13_compact_rdms,
    slide_14_snapshots,
    slide_12_temporal_fits,
    slide_15_model_diff,
    slide_16_temporal_mds,
    slide_19_workflow,
    slide_20_sliding_window,
    slide_17_discussion,
    slide_18_limitations,
]


def export_slide_previews(pptx_path: Path, out_dir: Path):
    """Export each slide as a JPEG via PowerPoint COM (Windows only).

    Falls back to a message if COM is unavailable.
    """
    out_dir.mkdir(parents=True, exist_ok=True)
    try:
        import comtypes.client
        import time
        ppt = comtypes.client.CreateObject("PowerPoint.Application")
        time.sleep(1)
        presentation = ppt.Presentations.Open(
            str(pptx_path.resolve()), WithWindow=False)
        for i, slide in enumerate(presentation.Slides, 1):
            slide.Export(
                str((out_dir / f"slide_{i:02d}.jpg").resolve()),
                "JPG", 1920, 1080)
        presentation.Close()
        try:
            ppt.Quit()
        except Exception:
            pass
        print(f"Exported {i} slide previews to {out_dir}")
    except ImportError:
        print("comtypes not available; skipping slide preview export. "
              "Run _export_slides.py separately on Windows.")


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    qc: dict = {"included_assets": [], "missing_assets": [],
                "required_coverage": {}, "slide_count": 0}
    assets = build_assets()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for idx, builder in enumerate(SLIDE_BUILDERS, 1):
        builder(prs, idx, assets, qc)

    qc["slide_count"] = len(prs.slides)
    qc["required_coverage"] = {
        "slide_count_10_to_20": 10 <= qc["slide_count"] <= 20,
        "change6_and_change5_static_rdms": (
            any("static_change6" in s for s in qc["included_assets"])
            and any("static_change5" in s for s in qc["included_assets"])),
        "change24_no1_temporal_rdms": (
            any("compact_row-temporal_change24_all-100-700ms-RT.png" in s
                for s in qc["included_assets"])
            and any("compact_row-temporal_change24_all-100-700ms-no1-RT.png" in s
                    for s in qc["included_assets"])),
        "temporal_fits_with_and_no1": (
            any("with_RT.png" in s and "temporal_model_fits" in s
                for s in qc["included_assets"])
            and any("no1-with_RT.png" in s and "temporal_model_fits" in s
                    for s in qc["included_assets"])),
        "mds_static_and_temporal": (
            any("mds_category-static_change6" in s
                for s in qc["included_assets"])
            and any("temporal_mds_static_0-700ms" in s
                    for s in qc["included_assets"])),
    }

    prs.save(str(OUTPUT_PPTX))
    OUTPUT_QC.write_text(json.dumps(qc, indent=2), encoding="utf-8")

    print(f"Created: {OUTPUT_PPTX}")
    print(f"QC: {OUTPUT_QC}")
    print(f"Slides: {qc['slide_count']}")
    print(f"Assets included: {len(qc['included_assets'])}")
    print(f"Assets missing: {len(qc['missing_assets'])}")

    # Export slide previews for visual QA
    export_slide_previews(OUTPUT_PPTX, PREVIEW_DIR)


if __name__ == "__main__":
    main()
