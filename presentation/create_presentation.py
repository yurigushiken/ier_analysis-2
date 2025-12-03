"""Programmatically build the Event Structure presentation with python-pptx.

The script favors a clean, earthy palette and injects the project-provided
media assets where possible. Run it from any working directory after
``conda activate ier_analysis`` to regenerate the PowerPoint artifact.
"""
from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
MEDIA_DIR = ROOT / "media"
OUTPUT = ROOT / "event_structure_infants.pptx"

EARTH_BG = RGBColor(245, 241, 235)
TITLE_COLOR = RGBColor(78, 52, 32)
BODY_COLOR = RGBColor(44, 62, 44)
ACCENT_GREEN = RGBColor(84, 120, 78)
ACCENT_RUST = RGBColor(170, 89, 57)
ACCENT_GOLD = RGBColor(199, 150, 75)

TITLE_SIZE = Pt(34)
BODY_SIZE = Pt(20)


def apply_background(slide) -> None:
    """Fill slides with the shared earthy tone."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = EARTH_BG


def style_title(title_shape) -> None:
    """Apply consistent styling to title placeholders."""
    text_frame = title_shape.text_frame
    text_frame.paragraphs[0].font.size = TITLE_SIZE
    text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    text_frame.paragraphs[0].font.bold = True


def configure_body_placeholder(body_shape) -> None:
    """Resize and style the main body placeholder for text content."""
    body_shape.left = Inches(0.5)
    body_shape.width = Inches(4.8)
    body_shape.top = Inches(1.5)
    body_shape.height = Inches(5.5)
    text_frame = body_shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = 0
    text_frame.margin_right = 0


def add_text_slide(prs: Presentation, title: str, bullets: Sequence[str]):
    """Add a slide with a title and bullet list."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_background(slide)
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    body_shape = slide.shapes.placeholders[1]
    configure_body_placeholder(body_shape)
    text_frame = body_shape.text_frame
    text_frame.clear()
    for idx, line in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.size = BODY_SIZE
        paragraph.font.color.rgb = BODY_COLOR
    return slide


def add_image(slide, filename: str, *, width: float = 4.7, top: float = 1.4) -> None:
    """Insert an image anchored to the right-hand side of the slide."""
    image_path = MEDIA_DIR / filename
    if not image_path.exists():
        raise FileNotFoundError(f"Media asset not found: {image_path}")
    slide.shapes.add_picture(
        str(image_path),
        Inches(5.4),
        Inches(top),
        width=Inches(width),
    )


def add_strategy_diagram(slide) -> None:
    """Render a simple diagram that visualizes the three gaze strategies."""
    shapes = slide.shapes
    nodes = {
        "Man Face": (Inches(0.8), Inches(4.3)),
        "Woman Face": (Inches(3.7), Inches(4.3)),
        "Body": (Inches(1.1), Inches(5.8)),
        "Toy": (Inches(3.4), Inches(5.8)),
        "Agent Face": (Inches(5.3), Inches(4.8)),
    }
    rendered = {}
    for label, (left, top) in nodes.items():
        shape = shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, left, top, Inches(1.3), Inches(0.7)
        )
        shape.fill.solid()
        fill_color = ACCENT_GREEN if "Face" in label else ACCENT_GOLD
        if label == "Toy":
            fill_color = ACCENT_RUST
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = BODY_COLOR
        shape.text = label
        shape.text_frame.paragraphs[0].font.size = Pt(14)
        shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        rendered[label] = shape

    def center(shape):
        return shape.left + shape.width / 2, shape.top + shape.height / 2

    def connect(start_label: str, end_label: str) -> None:
        start = center(rendered[start_label])
        end = center(rendered[end_label])
        connector = shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, start[0], start[1], end[0], end[1]
        )
        line = connector.line
        line.color.rgb = BODY_COLOR
        line.width = Pt(2)

    connect("Man Face", "Woman Face")
    connect("Woman Face", "Man Face")
    connect("Body", "Toy")
    connect("Toy", "Body")
    connect("Agent Face", "Toy")
    connect("Toy", "Agent Face")


def add_gw_vs_ugw_chart(slide) -> None:
    """Insert the GW vs UGW social verification paired bar chart."""
    chart_data = CategoryChartData()
    chart_data.categories = ["Agent-Agent Attention"]
    chart_data.add_series("GW (Upright)", (0.17,))
    chart_data.add_series("UGW (Inverted)", (0.008,))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(5.3),
        Inches(1.5),
        Inches(4.8),
        Inches(3.9),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.category_axis.tick_labels.font.size = Pt(14)
    chart.value_axis.maximum_scale = 0.2
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.tick_labels.number_format = "0%"
    chart.value_axis.tick_labels.font.size = Pt(12)

    for series, color in zip(chart.series, (ACCENT_GREEN, ACCENT_RUST)):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
        series.data_labels.show_value = True
        series.data_labels.number_format = "0.0%"
        series.data_labels.font.size = Pt(14)
        series.data_labels.font.bold = True


def build_presentation() -> None:
    """Assemble every slide and persist the PowerPoint file."""
    prs = Presentation()

    # Slide 1
    slide1 = add_text_slide(
        prs,
        "Definitions & Gaze Strategies",
        [
            "Concept: We moved beyond simple \"looking time\" to analyze every AOI transition and classify each gaze strategy.",
            "Strategy 1 - Man Face ↔ Woman Face (Agent-Agent Attention): Direct gaze shifts between the giver and receiver that verify shared intent.",
            "Strategy 2 - Body ↔ Toy (Motion Tracking): Tracks the hands/bodies and the moving object without invoking the agents' minds.",
            "Strategy 3 - Face ↔ Toy (Agent-Object Binding): Links an agent's face and the object yet stops short of connecting both people.",
            "Visual: Diagram shows every bidirectional arrow for the three strategies.",
        ],
    )
    add_strategy_diagram(slide1)

    # Slide 2
    slide2 = add_text_slide(
        prs,
        "The Rise of Agent-Agent Attention",
        [
            "Story: Older infants increasingly verify the social connection between actors.",
            "Key Finding: Face ↔ Face transitions climb steadily from 7 to 11 months (~4 percentage points per month).",
            "Stats: Linear trend across 7-11 months is significant (p = 0.025).",
            "Adults Comparison: Adults show far higher Agent-Agent Attention than 7-month-olds (p = 0.006).",
            "Visual: Linear trend scatter with regression line.",
        ],
    )
    add_image(slide2, "gw_transitions_min3_50_percent_linear_trend_agent_agent_attention.png")

    # Slide 3
    slide3 = add_text_slide(
        prs,
        "The Decline of Motion Tracking",
        [
            "Story: As Agent-Agent Attention rises, purely motion tracking falls.",
            "Key Finding: 7-month-olds devote over half of their transitions to Body ↔ Toy tracking before the strategy drops with age.",
            "Stats: Linear trend shows a significant negative slope (p = 0.005).",
            "Baseline: Motion Tracking dominates early gaze behavior.",
            "Visual: Linear trend scatter for Motion Tracking.",
        ],
    )
    add_image(
        slide3,
        "gw_transitions_min3_50_percent_linear_trend_motion_tracking.png",
    )

    # Slide 4
    slide4 = add_text_slide(
        prs,
        "Agent-Object Binding",
        [
            "Story: Face ↔ Toy transitions peak at 9 months, bridging motion tracking and full agent-agent attention.",
            "Observation: 9-month-olds allocate roughly half of their transitions to linking the agent and the object but not the recipient.",
            "Stats: Placeholder for the detailed GEE stats describing the 9-month peak and adult contrasts.",
            "Visual: Developmental profile of Agent-Object Binding.",
        ],
    )
    add_image(
        slide4,
        "gw_transitions_min3_50_percent_linear_trend_agent_object_binding.png",
    )

    # Slide 5
    slide5 = add_text_slide(
        prs,
        "Semantic Control (GW vs UGW)",
        [
            "Question: Do 10-month-olds fixate faces because of semantic understanding or simple facial attraction?",
            "Method: Compare the upright Give (GW) to the inverted Give (UGW).",
            "Key Finding: GW maintains high Agent-Agent Attention (17%) but UGW collapses to 0.8%.",
            "Conclusion: Strategy is driven by semantics; the gaze pattern disappears when meaning breaks.",
            "Visual: Paired bar chart (GW vs UGW) for Agent-Agent Attention.",
        ],
    )
    add_gw_vs_ugw_chart(slide5)

    # Slide 6
    slide6 = add_text_slide(
        prs,
        "The Anticipation Gap (Latency to Toy)",
        [
            "Concept: Theory of Mind should enable predictive saccades toward the toy before the transfer completes.",
            "Finding: Adults anticipate instantly (4.0 frames) while infants remain reactive at 13-22 frames.",
            "Stats: Every infant cohort is significantly slower than adults (p < 0.01).",
            "Conclusion: Eleven-month-olds understand structure (faces) but still process rather than predict.",
            "Visual: Mean latency bar chart with significance markers.",
        ],
    )
    add_image(slide6, "gw_latency_to_toy_min3_50_percent_latency_to_toy_bar.png")

    # Slide 7
    slide7 = add_text_slide(
        prs,
        "The Reaction Check (Show Condition)",
        [
            "Concept: In the Show event, success requires monitoring the observer's (Man's) reaction.",
            "Finding: Adults check the Man's face 52% of the time; 7-month-olds never check and 11-month-olds reach only 9%.",
            "Stats: All infant cohorts are significantly lower than adults (p < 0.001).",
            "Conclusion: Pre-linguistic event representation is partial: agent and object are encoded, but the social goal is missed.",
            "Visual: Reaction-window bar chart for looking at the Man.",
        ],
    )
    add_image(slide7, "sw_time_window_min3_50_percent_time_window_bar_plot.png")

    # Slide 8
    add_text_slide(
        prs,
        "The Story in Summary",
        [
            "7-8 Months: Motion tracking driven by Body ↔ Toy transitions with reflexively fast reactions.",
            "9 Months: Agent-object binders obsessed with Agent ↔ Object connections and slower processing times.",
            "10-11 Months: Agent-agent attention peaks—infants capture giver, recipient, and object but still process slowly and miss reactions.",
            "Adults: Experts who are fast, predictive, and fully socially integrated.",
            "Takeaway: Event structure understanding emerges in stages; efficiency and social goals lag behind structural mastery.",
        ],
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Presentation saved to {OUTPUT}")


def main() -> None:
    build_presentation()


if __name__ == "__main__":
    main()
